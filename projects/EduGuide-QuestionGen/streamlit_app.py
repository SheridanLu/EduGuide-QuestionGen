# streamlit_app.py - 美观简约版 v4.1
import streamlit as st
import json
import os

# ========== 多语言 ==========
LANGUAGES = {"en": "🇺🇸 EN", "zh-CN": "🇨🇳 中文", "zh-TW": "🇹🇼 繁體"}

T = {
    "en": {
        "title": "EduGuide", "subtitle": "AI-Powered Question Generation",
        "upload": "Upload Material", "upload_hint": "Drag & drop or click to upload",
        "input": "Or paste text...", "words": "words",
        "error": "Student error (optional)", "error_hint": "e.g., confused X with Y",
        "generate": "Generate", "generating": "Generating...", "done": "Done!",
        "knowledge": "Knowledge", "questions": "Questions", "remedial": "Remedial",
        "basic": "Basic", "intermediate": "Intermediate", "advanced": "Advanced",
        "api": "Settings", "provider": "Provider", "key": "API Key", "model": "Model",
        "save": "Save", "saved": "Saved!", "language": "Language",
        "start": "Start Practice", "step": "Step", "answer": "Your Answer",
        "submit": "Submit", "hint": "Hint", "correct": "Correct!",
        "retry": "Try again", "complete": "Complete",
        "progress": "Progress", "back": "Back",
        "current_api": "Current API", "files_loaded": "files loaded",
        "clear": "Clear",
    },
    "zh-CN": {
        "title": "EduGuide", "subtitle": "智能出题系统",
        "upload": "上传教材", "upload_hint": "拖拽或点击上传",
        "input": "或粘贴文本...", "words": "字",
        "error": "学生错误（可选）", "error_hint": "例如：混淆了A和B",
        "generate": "生成题目", "generating": "生成中...", "done": "完成！",
        "knowledge": "知识点", "questions": "题目", "remedial": "补救",
        "basic": "基础", "intermediate": "中级", "advanced": "高级",
        "api": "API 设置", "provider": "提供商", "key": "密钥", "model": "模型",
        "save": "保存", "saved": "已保存！", "language": "语言",
        "start": "开始解题", "step": "步骤", "answer": "你的答案",
        "submit": "提交", "hint": "提示", "correct": "正确！", "retry": "再试一次",
        "complete": "完成", "progress": "进度", "back": "返回",
        "current_api": "当前 API", "files_loaded": "个文件已加载",
        "clear": "清空",
    },
    "zh-TW": {
        "title": "EduGuide", "subtitle": "智慧出題系統",
        "upload": "上傳教材", "upload_hint": "拖曳或點擊上傳",
        "input": "或貼上文字...", "words": "字",
        "error": "學生錯誤（可選）", "error_hint": "例如：混淆了A和B",
        "generate": "生成題目", "generating": "生成中...", "done": "完成！",
        "knowledge": "知識點", "questions": "題目", "remedial": "補救",
        "basic": "基礎", "intermediate": "中級", "advanced": "高級",
        "api": "API 設定", "provider": "提供商", "key": "金鑰", "model": "模型",
        "save": "儲存", "saved": "已儲存！", "language": "語言",
        "start": "開始解題", "step": "步驟", "answer": "你的答案",
        "submit": "提交", "hint": "提示", "correct": "正確！", "retry": "再試一次",
        "complete": "完成", "progress": "進度", "back": "返回",
        "current_api": "當前 API", "files_loaded": "個檔案已載入",
        "clear": "清空",
    }
}

def t(key, lang="en"):
    return T.get(lang, T["en"]).get(key, key)

# ========== 配置 ==========
st.set_page_config(page_title="EduGuide", page_icon="🎓", layout="wide", initial_sidebar_state="auto")

# ========== CSS ==========
st.markdown("""<style>
/* === Reset === */
#MainMenu, footer, header {visibility: hidden}
.css-1d391kg, .ezpmspp {display: none}

/* === Global === */
.main .block-container {
    padding-top: 0;
    padding-bottom: 4rem;
    max-width: 960px;
}

* {
    font-family: -apple-system, BlinkMacSystemFont, "SF Pro Display", "Segoe UI", Roboto, sans-serif !important;
}

/* === Top Navigation === */
.top-nav {
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 20px 0 16px 0;
    margin-bottom: 0;
}

.logo {
    font-size: 1.35rem;
    font-weight: 700;
    color: #111;
    letter-spacing: -0.5px;
    display: flex;
    align-items: center;
    gap: 8px;
}

.logo-icon {
    width: 36px;
    height: 36px;
    background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%);
    border-radius: 10px;
    display: flex;
    align-items: center;
    justify-content: center;
    color: white;
    font-size: 18px;
}

/* === Language Buttons (segmented) === */
.lang-col > div {
    padding: 0 !important;
}
.lang-col button {
    border-radius: 10px !important;
    font-weight: 600 !important;
    font-size: 0.88rem !important;
    padding: 9px 0 !important;
    border: 1.5px solid #e8e8e8 !important;
    background: white !important;
    color: #888 !important;
    transition: all 0.2s ease !important;
    box-shadow: none !important;
}
.lang-col button:hover:not(:disabled) {
    border-color: #c4b5fd !important;
    color: #6366f1 !important;
    background: #faf8ff !important;
}
.lang-col button:disabled {
    background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%) !important;
    color: white !important;
    border-color: transparent !important;
    box-shadow: 0 2px 8px rgba(99, 102, 241, 0.3) !important;
}
.lang-col button:disabled:hover {
    background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%) !important;
    color: white !important;
}
.api-col > div {
    padding: 0 !important;
}
.api-col button {
    border-radius: 10px !important;
    font-size: 1.1rem !important;
    padding: 9px 0 !important;
    border: 1.5px solid #e8e8e8 !important;
    background: white !important;
    color: #aaa !important;
    transition: all 0.2s ease !important;
    box-shadow: none !important;
}
.api-col button:hover {
    border-color: #c4b5fd !important;
    color: #6366f1 !important;
}

/* === Hero Section === */
.hero-section {
    text-align: center;
    padding: 48px 24px 40px;
    margin-bottom: 8px;
}

.hero-badge {
    display: inline-block;
    padding: 6px 16px;
    background: #f0f0ff;
    color: #6366f1;
    border-radius: 100px;
    font-size: 0.8rem;
    font-weight: 600;
    margin-bottom: 20px;
    letter-spacing: 0.3px;
}

.hero-title {
    font-size: 2.8rem;
    font-weight: 800;
    color: #111;
    margin: 0 0 12px 0;
    letter-spacing: -1.5px;
    line-height: 1.1;
}

.hero-desc {
    font-size: 1.1rem;
    color: #666;
    margin: 0;
    font-weight: 400;
    max-width: 480px;
    margin: 0 auto;
    line-height: 1.6;
}

/* === Card === */
.glass-card {
    background: #fff;
    border-radius: 20px;
    padding: 28px;
    border: 1px solid #f0f0f0;
    box-shadow: 0 1px 3px rgba(0,0,0,0.03), 0 4px 16px rgba(0,0,0,0.02);
    margin-bottom: 12px;
}

/* === Upload Area === */
.upload-zone {
    border: 2px dashed #e0e0e0;
    border-radius: 16px;
    padding: 36px 24px;
    text-align: center;
    background: #fafbfc;
    transition: all 0.3s ease;
    cursor: pointer;
    margin-bottom: 20px;
}

.upload-zone:hover {
    border-color: #6366f1;
    background: #f8f8ff;
}

.upload-icon {
    font-size: 2rem;
    margin-bottom: 10px;
    opacity: 0.6;
}

.upload-label {
    color: #888;
    font-size: 0.9rem;
    font-weight: 500;
}

.upload-formats {
    color: #bbb;
    font-size: 0.78rem;
    margin-top: 6px;
}

/* === Input === */
.stTextArea textarea {
    border: 1.5px solid #eee !important;
    border-radius: 14px !important;
    padding: 16px 18px !important;
    font-size: 0.95rem !important;
    line-height: 1.7 !important;
    transition: border-color 0.2s !important;
    background: #fafbfc !important;
}

.stTextArea textarea:focus {
    border-color: #6366f1 !important;
    background: #fff !important;
    box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.08) !important;
}

.stTextInput > div > div {
    border-radius: 12px !important;
    border: 1.5px solid #eee !important;
    padding: 8px 14px !important;
}

.stTextInput > div > div:focus-within {
    border-color: #6366f1 !important;
    box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.08) !important;
}

/* === Select === */
.stSelectbox > div > div {
    border-radius: 10px !important;
    border: 1.5px solid #eee !important;
    padding: 6px 12px !important;
}

/* === Buttons === */
.stButton > button {
    border-radius: 12px !important;
    font-weight: 600 !important;
    font-size: 0.92rem !important;
    padding: 10px 24px !important;
    transition: all 0.2s ease !important;
    border: none !important;
}

.stButton > button:hover {
    transform: translateY(-1px);
    box-shadow: 0 4px 14px rgba(0,0,0,0.1);
}

.stButton > button[data-testid="stPrimary"],
.stButton > button[kind="primary"] {
    background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%) !important;
    color: white !important;
    box-shadow: 0 2px 10px rgba(99, 102, 241, 0.3) !important;
}

.stButton > button[data-testid="stPrimary"]:hover,
.stButton > button[kind="primary"]:hover {
    box-shadow: 0 4px 18px rgba(99, 102, 241, 0.4) !important;
    background: linear-gradient(135deg, #5558e6 0%, #7c4fe0 100%) !important;
}

/* === Generate Button (large) === */
.generate-btn {
    text-align: center;
    padding: 16px 0;
}

/* === API Card === */
.api-badge {
    display: inline-flex;
    align-items: center;
    gap: 8px;
    padding: 10px 18px;
    background: #f8f8ff;
    border: 1px solid #ede9fe;
    border-radius: 12px;
    font-size: 0.88rem;
    color: #6366f1;
    font-weight: 600;
}

.api-badge .api-dot {
    width: 8px;
    height: 8px;
    background: #22c55e;
    border-radius: 50%;
    animation: pulse 2s infinite;
}

@keyframes pulse {
    0%, 100% { opacity: 1; }
    50% { opacity: 0.4; }
}

/* === Section Title === */
.section-title {
    font-size: 0.8rem;
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: 1px;
    color: #aaa;
    margin-bottom: 12px;
    margin-top: 20px;
}

/* === Tabs === */
.stTabs [data-baseweb="tab-list"] {
    gap: 4px;
    background: transparent;
    border-bottom: 1px solid #f0f0f0;
    margin-bottom: 20px;
}

.stTabs [data-baseweb="tab"] {
    border-radius: 0 !important;
    padding: 12px 24px !important;
    font-weight: 500 !important;
    font-size: 0.9rem !important;
    background: transparent !important;
    color: #aaa !important;
    border-bottom: 2px solid transparent !important;
    transition: all 0.2s ease !important;
}

.stTabs [aria-selected="true"] {
    color: #111 !important;
    border-bottom: 2px solid #6366f1 !important;
    font-weight: 600 !important;
}

/* === Question Card === */
.question-card {
    background: #fafbfc;
    border-radius: 14px;
    padding: 18px 22px;
    margin-bottom: 10px;
    border: 1px solid #f0f0f0;
    transition: all 0.2s;
}

.question-card:hover {
    border-color: #ddd;
    box-shadow: 0 2px 8px rgba(0,0,0,0.04);
}

/* === Step Cards === */
.step-done {
    background: linear-gradient(135deg, #f0fdf4 0%, #dcfce7 100%);
    border-left: 3px solid #22c55e;
    border-radius: 14px;
    padding: 16px 20px;
    margin: 8px 0;
}

.step-current {
    background: #fff;
    border: 1.5px solid #ede9fe;
    border-radius: 14px;
    padding: 20px;
    margin: 8px 0;
    box-shadow: 0 2px 12px rgba(99, 102, 241, 0.06);
}

/* === Divider === */
.divider {
    height: 1px;
    background: linear-gradient(90deg, transparent 0%, #eee 50%, transparent 100%);
    border: none;
    margin: 24px 0;
}

/* === Misc === */
.small { font-size: 0.82rem; color: #bbb; }
.meta { color: #999; font-size: 0.85rem; }
</style>
""", unsafe_allow_html=True)

# ========== 初始化 ==========
if 'lang' not in st.session_state: st.session_state.lang = "zh-CN"
if 'text' not in st.session_state: st.session_state.text = ""
if 'result' not in st.session_state: st.session_state.result = None
if 'practice' not in st.session_state: st.session_state.practice = None
if 'step' not in st.session_state: st.session_state.step = 0
if 'answers' not in st.session_state: st.session_state.answers = {}
if 'show_api' not in st.session_state: st.session_state.show_api = False

lang = st.session_state.lang

# ========== 顶部导航 ==========
st.markdown(f"""
<div class="top-nav">
    <div class="logo">
        <div class="logo-icon">🎓</div>
        EduGuide
    </div>
    <div class="nav-actions">
        <span class="api-badge">
            <span class="api-dot"></span>
            AI Ready
        </span>
    </div>
</div>
""", unsafe_allow_html=True)

# 语言切换按钮（Streamlit原生）
col1, col2, col3, col4 = st.columns([1, 1, 1, 0.6])
for i, (code, label) in enumerate(LANGUAGES.items()):
    with [col1, col2, col3][i]:
        st.markdown('<div class="lang-col">', unsafe_allow_html=True)
        if code == lang:
            st.button(label, disabled=True, use_container_width=True)
        else:
            if st.button(label, use_container_width=True):
                st.session_state.lang = code
                st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
with col4:
    st.markdown('<div class="api-col">', unsafe_allow_html=True)
    if st.button(f"⚙️", key="api_toggle", use_container_width=True):
        st.session_state.show_api = not st.session_state.show_api
    st.markdown('</div>', unsafe_allow_html=True)

# ========== Hero ==========
st.markdown(f"""
<div class="hero-section">
    <div class="hero-badge">AI-Powered Education</div>
    <div class="hero-title">{t('title', lang)}</div>
    <div class="hero-desc">{t('subtitle', lang)}</div>
</div>
""", unsafe_allow_html=True)

# ========== API设置 ==========
if st.session_state.show_api:
    st.markdown("---")
    try:
        from config.api_config import get_config_manager, APIProvider, ProviderConfig
        cm = get_config_manager()
        current = cm.current_provider
        providers = cm.get_all_providers()
        options = {p['id']: p['name'] for p in providers}
        
        selected = st.selectbox(
            f"**{t('provider', lang)}**",
            list(options.keys()),
            format_func=lambda x: options[x],
            index=list(options.keys()).index(current.value),
        )
        config = cm.providers.get(APIProvider(selected))
        
        col1, col2 = st.columns(2)
        with col1:
            key = st.text_input(f"**{t('key', lang)}**", value=config.api_key, type="password")
        with col2:
            model = st.text_input(f"**{t('model', lang)}**", value=config.model)
            if st.button(f"✅ {t('save', lang)}", type="primary", use_container_width=True):
                new_cfg = ProviderConfig(name=config.name, api_key=key, base_url=config.base_url, model=model, enabled=True)
                cm.update_provider_config(APIProvider(selected), new_cfg)
                cm.set_current_provider(APIProvider(selected))
                st.success(t('saved', lang))
    except Exception as e:
        st.error(f"Error: {e}")
    st.markdown("---")

# ========== 主内容 ==========
st.markdown('<div class="section-title">INPUT</div>', unsafe_allow_html=True)
st.markdown(f'<div class="glass-card">', unsafe_allow_html=True)

# 上传区域
st.markdown(f"""
<div class="upload-zone">
    <div class="upload-icon">📄</div>
    <div class="upload-label">{t('upload_hint', lang)}</div>
    <div class="upload-formats">TXT · PDF · DOCX · PPTX</div>
</div>
""", unsafe_allow_html=True)

uploaded_files = st.file_uploader(
    t('upload', lang),
    type=['txt', 'pdf', 'docx', 'pptx'],
    accept_multiple_files=True,
    label_visibility="collapsed"
)
if uploaded_files:
    combined = []
    file_names = []
    for f in uploaded_files:
        try:
            f.seek(0)
            if f.name.endswith('.txt'):
                combined.append(f.read().decode('utf-8'))
            elif f.name.endswith('.pdf'):
                import PyPDF2
                pages = PyPDF2.PdfReader(f).pages
                combined.append("\n".join([p.extract_text() for p in pages if p.extract_text()]))
            elif f.name.endswith('.docx'):
                import docx
                combined.append("\n".join([p.text for p in docx.Document(f).paragraphs if p.text]))
            elif f.name.endswith('.pptx'):
                from pptx import Presentation
                text = ""
                for slide in Presentation(f).slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"): text += shape.text + "\n"
                combined.append(text)
            file_names.append(f.name)
        except Exception as e:
            st.error(f"Error reading {f.name}: {e}")
    if combined:
        st.session_state.text = "\n\n---\n\n".join(combined)
        st.success(f"✅ {len(file_names)} {t('files_loaded', lang)}: {', '.join(file_names)}")

# 文本输入
text = st.text_area(
    t('input', lang),
    value=st.session_state.text,
    height=180,
    placeholder=t('placeholder', lang) if 'placeholder' in T.get(lang, {}) else t('input', lang),
    label_visibility="collapsed"
)
if text:
    st.session_state.text = text
    st.markdown(f"<div class='small' style='text-align:right'>{len(text)} {t('words', lang)}</div>", unsafe_allow_html=True)

st.markdown('</div>', unsafe_allow_html=True)  # close glass-card

# 学生错误点
col_err, _ = st.columns([3, 2])
with col_err:
    error = st.text_input(
        f"💡 {t('error', lang)}",
        placeholder=t('error_hint', lang),
        label_visibility="collapsed"
    )

# 生成按钮
st.markdown('<div class="divider"></div>', unsafe_allow_html=True)
st.markdown('<div class="generate-btn">', unsafe_allow_html=True)
if st.button(f"🚀  {t('generate', lang)}", type="primary"):
    if not text:
        st.warning("Please input" if lang == "en" else "请输入内容")
    else:
        with st.spinner(t('generating', lang)):
            try:
                from workflow.openclaw_flow import OpenClawFlow
                st.session_state.result = OpenClawFlow().run(text, error if error else None)
                st.success(t('done', lang))
            except Exception as e:
                st.error(f"Error: {e}")
st.markdown('</div>', unsafe_allow_html=True)

# ========== 结果 ==========
if st.session_state.result:
    st.markdown('<div class="divider"></div>', unsafe_allow_html=True)
    st.markdown('<div class="section-title">RESULTS</div>', unsafe_allow_html=True)
    
    q_data, a_data = {}, {}
    if os.path.exists("output/questions.json"):
        with open("output/questions.json", 'r', encoding='utf-8') as f:
            q_data = json.load(f)
    if os.path.exists("output/answers.json"):
        with open("output/answers.json", 'r', encoding='utf-8') as f:
            a_data = json.load(f)
    
    if st.session_state.practice:
        qid = st.session_state.practice
        level, idx = qid.split('_')
        idx = int(idx)
        question = q_data.get(level, [])[idx]
        guidance = a_data.get(level, [])[idx].get('guidance', {}) if idx < len(a_data.get(level, [])) else {}
        steps = [guidance.get(f'step{i}', '') for i in range(1, 4)]
        steps = [s for s in steps if s]
        
        if st.button(f"← {t('back', lang)}"):
            st.session_state.practice = None
            st.session_state.step = 0
            st.rerun()
        
        st.markdown(f"### {question}")
        progress = st.session_state.step / len(steps) if steps else 0
        st.progress(progress)
        st.markdown(f"**{t('progress', lang)}: {st.session_state.step}/{len(steps)}**")
        
        for i in range(st.session_state.step):
            ans = st.session_state.answers.get(f"{qid}_{i}", '')
            st.markdown(f"<div class='step-done'><strong>✅ {t('step', lang)} {i+1}:</strong> {steps[i]}<br><small class='meta'>→ {ans}</small></div>", unsafe_allow_html=True)
        
        if st.session_state.step < len(steps):
            st.markdown(f"<div class='step-current'><strong>{t('step', lang)} {st.session_state.step+1}:</strong> {steps[st.session_state.step]}</div>", unsafe_allow_html=True)
            ans = st.text_area(t('answer', lang), key=f"ans_{qid}_{st.session_state.step}", height=100)
            c1, c2 = st.columns(2)
            with c1:
                if st.button(f"✅ {t('submit', lang)}", use_container_width=True):
                    if ans.strip():
                        st.session_state.answers[f"{qid}_{st.session_state.step}"] = ans
                        st.session_state.step += 1
                        st.success(t('correct', lang))
                        st.rerun()
                    else:
                        st.warning(t('retry', lang))
            with c2:
                if st.button(f"💡 {t('hint', lang)}", use_container_width=True):
                    hints = guidance.get('hints', [])
                    st.info(hints[st.session_state.step] if st.session_state.step < len(hints) else "Think about key concepts.")
        else:
            st.success(t('complete', lang))
            if st.button(f"✅ {t('complete', lang)}", use_container_width=True):
                st.session_state.practice = None
                st.session_state.step = 0
                st.rerun()
    else:
        t1, t2, t3 = st.tabs([f"💡 {t('knowledge', lang)}", f"📝 {t('questions', lang)}", f"🤝 {t('remedial', lang)}"])
        
        with t1:
            if os.path.exists("output/knowledge.json"):
                with open("output/knowledge.json", 'r', encoding='utf-8') as f:
                    for i, p in enumerate(json.load(f).get("knowledge_points", []), 1):
                        st.markdown(f"**{i}.** {p}")
        
        with t2:
            for level, title in [("basic", f"🟢 {t('basic', lang)}"), ("intermediate", f"🟡 {t('intermediate', lang)}"), ("advanced", f"🔴 {t('advanced', lang)}")]:
                if level in q_data:
                    st.markdown(f"### {title}")
                    for i, q in enumerate(q_data[level], 1):
                        st.markdown(f"<div class='question-card'><strong>Q{i}:</strong> {q}</div>", unsafe_allow_html=True)
                        if st.button(t('start', lang), key=f"start_{level}_{i}", type="primary"):
                            st.session_state.practice = f"{level}_{i}"
                            st.session_state.step = 0
                            st.rerun()
        
        with t3:
            if os.path.exists("output/remedial.json"):
                with open("output/remedial.json", 'r', encoding='utf-8') as f:
                    for item in json.load(f).get("remedial", []):
                        g = item.get('guidance', {})
                        for k in ['probing_question_1', 'probing_question_2', 'probing_question_3']:
                            if k in g:
                                st.markdown(f"**{k.replace('_', ' ').title()}:** {g[k]}")
