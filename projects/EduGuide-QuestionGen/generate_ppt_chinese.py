# generate_ppt_chinese.py - 生成中文版PPT v2.0
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色
PURPLE = RGBColor(99, 102, 241)
PURPLE_DARK = RGBColor(139, 92, 246)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(26, 26, 26)
GRAY = RGBColor(128, 128, 128)
LIGHT_BG = RGBColor(245, 245, 255)
GREEN = RGBColor(34, 197, 94)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)

def set_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # 移到底层
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_bullet(slide, left, top, width, height, items, size=16, color=BLACK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf

def add_title_bar(slide, title_text, subtitle_text=""):
    set_bg(slide)
    # 紫色标题栏
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    add_text(slide, 0.8, 0.2, 11, 0.8, title_text, 32, True, WHITE)
    if subtitle_text:
        add_text(slide, 0.8, 0.75, 11, 0.4, subtitle_text, 14, False, RGBColor(200, 200, 255))

def add_footer(slide, page_num):
    add_text(slide, 11.5, 7.1, 1.5, 0.3, f"{page_num}/25", 11, False, GRAY, PP_ALIGN.RIGHT)
    add_text(slide, 0.5, 7.1, 4, 0.3, "EduGuide - INT2094 NLP", 9, False, GRAY)

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PURPLE)
add_text(slide, 1, 2, 11, 1.5, "EduGuide - 智能出题系统", 48, True, WHITE, PP_ALIGN.CENTER)
add_text(slide, 1, 3.5, 11, 0.8, "基于苏格拉底式引导的AI驱动教育解决方案", 22, False, RGBColor(220, 220, 255), PP_ALIGN.CENTER)
add_text(slide, 1, 5, 11, 1.5, "课程：INT2094 - 自然语言处理\n小组：[小组标签] | 日期：[演讲日期]", 16, False, RGBColor(200, 200, 255), PP_ALIGN.CENTER)

# ========== 第2页：议程 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "演讲议程")
items = [
    "1. 📖 引言 — 我们要解决的问题",
    "2. 🏗️ 系统设计 — 多Agent架构与技术栈",
    "3. 💻 实现与演示 — 功能展示与现场Demo",
    "4. 📊 评估 — 优势、劣势与对比分析",
    "5. 🎯 结论 — 总结与未来工作",
    "6. ❓ 问答 — Q&A",
]
add_bullet(slide, 1.5, 1.8, 10, 5, items, 22)
add_footer(slide, 2)

# ========== 第3页：什么是EduGuide ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "EduGuide：智能出题系统")
add_text(slide, 1, 1.5, 11, 0.8, "一个基于AI的系统，从教学材料自动生成题目，采用苏格拉底式引导，并通过AI实时验证学生答案。", 18, False, GRAY)
items = [
    "📚 面向教师：节省出题时间（数小时 → 数秒）",
    "🎯 面向学生：引导式发现学习，而非直接给答案",
    "🌍 面向所有人：多语言支持（英/中/繁）",
]
add_bullet(slide, 1, 2.8, 5.5, 3, items, 18)
items2 = [
    "✅ 自动题目生成（分层难度）",
    "✅ 苏格拉底式分步引导",
    "✅ AI答案验证 + 错因分析",
    "✅ 交互式练习（无限次重试）",
    "✅ 多文件同时上传",
    "✅ Agent监控可视化",
    "✅ 多语言界面",
]
add_bullet(slide, 7, 2.8, 5.5, 3.5, items2, 16)
add_footer(slide, 3)

# ========== 第4页：挑战 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "挑战：创建优质题目很困难")
add_text(slide, 1, 1.5, 5.5, 0.5, "教师面临的挑战", 22, True, PURPLE)
items1 = ["⏰ 创建有深度的题目需要数小时", "🎯 难以创建不同难度的题目", "📝 难以提供个性化指导", "🔄 每次新材料都要重复过程"]
add_bullet(slide, 1, 2.2, 5.5, 2.5, items1, 16)
add_text(slide, 7, 1.5, 5.5, 0.5, "学生面临的挑战", 22, True, PURPLE)
items2 = ["🤔 直接答案不促进深度学习", "📉 缺乏分步指导", "❌ 犯错后无反馈", "🔄 不知道为什么错、如何纠正"]
add_bullet(slide, 7, 2.2, 5.5, 2.5, items2, 16)
add_footer(slide, 4)

# ========== 第5页：现有方案局限 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "现有方案：好，但还不够")
items = [
    "ChatGPT：生成题目 → ❌ 直接给答案、无分步验证、无错因分析",
    "Quizlet：抽认卡 → ❌ 需手动创建、无AI生成",
    "Khan Academy：练习题 → ❌ 仅预制内容、非从材料生成",
    "Google Forms：测验 → ❌ 需手动出题、无AI辅助",
]
add_bullet(slide, 1, 1.8, 11, 2.5, items, 17)
add_text(slide, 1, 5, 11, 1, "❌ 没有任何方案能同时做到：从材料生成题目 + 苏格拉底式引导 + AI答案验证 + 交互式练习", 18, True, RED)
add_footer(slide, 5)

# ========== 第6页：我们的方案 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "EduGuide：填补空白")
items = [
    "✅ 自动生成 — 从材料生成题目（数小时→数秒）",
    "✅ 苏格拉底式引导 — 分步引导，不直接告诉",
    "✅ AI答案验证 — LLM判断答案正确性",
    "✅ 错因分析 — 答错时分析原因+引导",
    "✅ 交互式练习 — 无限次重试",
    "✅ 分层题目 — 基础/中级/高级",
    "✅ 多文件上传 — 同时上传多个文件",
    "✅ 多语言 — 英/中/繁",
]
add_bullet(slide, 1, 1.6, 11, 5, items, 17)
add_footer(slide, 6)

# ========== 第7页：Agent与LLM的关系 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Agent与LLM：它们是什么关系？")
add_text(slide, 1, 1.5, 5.5, 0.6, "什么是LLM（大语言模型）？", 18, True, PURPLE)
add_bullet(slide, 1, 2.2, 5.5, 1.2, ["AI的大脑 — 能理解和生成人类语言", "例子：GPT-4, GLM-4, DeepSeek, Claude", "像一个知识渊博但需要明确指令的专家"], 14)
add_text(slide, 7, 1.5, 5.5, 0.6, "什么是Agent（智能代理）？", 18, True, PURPLE)
add_bullet(slide, 7, 2.2, 5.5, 1.2, ["有特定任务的\"虚拟工作者\" — 调用LLM完成具体工作", "类比：LLM=万能工具 | Agent=使用工具的工匠", "比直接用LLM更有结构和目的性"], 14)
add_text(slide, 1, 3.8, 5.5, 0.4, "为什么不直接用LLM？", 16, True, RED)
add_bullet(slide, 1, 4.4, 5.5, 1.5, [
    "❌ 直接用LLM：一次要求太多 → 输出混乱、质量不稳定",
    "✅ 用Agent：每个只做一件事 → 质量高、可调试、可监控",
    "就像：一个人做所有事 vs. 专家团队各司其职",
], 14)
add_text(slide, 7, 3.8, 5.5, 0.4, "EduGuide中的Agent分工", 16, True, GREEN)
add_bullet(slide, 7, 4.4, 5.5, 2, [
    "📚 知识Agent → \"提取关键概念\"",
    "📝 题目Agent → \"生成分层题目\"",
    "🎯 答案Agent → \"设计苏格拉底引导\"",
    "🤝 补救Agent → \"生成补救建议\"",
], 14)
add_footer(slide, 26)

# ========== 第8页：系统架构 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "多Agent系统架构")
flow = [
    "📁 教学材料输入 (文本/PDF/DOCX/PPTX，支持多文件)",
    "        ↓",
    "📚 Agent 1: 知识提取 → knowledge.json (~3s)",
    "        ↓",
    "📝 Agent 2: 题目生成 → questions.json (~4s)",
    "        ↓",
    "🎯 Agent 3: 答案引导 → answers.json (~6s)",
    "        ↓",
    "🤝 Agent 4: 补救支持 → remedial.json (~3s)",
    "        ↓",
    "🎓 交互式练习界面 (AI验证 + 错因分析 + 无限重试)",
]
add_bullet(slide, 2, 1.6, 9, 5.5, flow, 16)
add_footer(slide, 26)

# ========== 第8页：多Agent设计 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "为什么采用多Agent架构？")
items = [
    "📚 知识Agent — 提取关键概念 (~3s)",
    "📝 题目Agent — 生成分层题目 (~4s)",
    "🎯 答案Agent — 创建苏格拉底引导 (~6s)",
    "🤝 补救Agent — 解决学生错误 (~3s)",
]
add_bullet(slide, 1, 1.6, 6, 3, items, 17)
add_text(slide, 7.5, 1.6, 5, 0.5, "设计优势", 22, True, PURPLE)
items2 = ["✅ 模块化 — 独立更新各个Agent", "✅ 可调试 — Agent Monitor实时监控", "✅ 可扩展 — 为新功能添加新Agent", "✅ 可维护 — 每个Agent只做一件事"]
add_bullet(slide, 7.5, 2.3, 5, 3, items2, 16)
add_footer(slide, 26)

# ========== 第9页：技术栈 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "技术栈")
add_text(slide, 1, 1.5, 4, 0.5, "前端", 20, True, PURPLE)
add_bullet(slide, 1, 2.2, 4, 1.5, ["Streamlit (Python)", "自定义CSS — 紫色渐变主题"], 15)
add_text(slide, 5.5, 1.5, 4, 0.5, "后端", 20, True, PURPLE)
add_bullet(slide, 5.5, 2.2, 4, 1.5, ["Python 3.12", "JSON文件通信"], 15)
add_text(slide, 10, 1.5, 3, 0.5, "文件处理", 20, True, PURPLE)
add_bullet(slide, 10, 2.2, 3, 1.5, ["PyPDF2", "python-docx", "python-pptx"], 15)
add_text(slide, 1, 4, 11, 0.5, "AI/LLM提供商（7个）", 20, True, PURPLE)
items = ["🟣 智谱GLM (glm-4-plus) | 🔵 DeepSeek (deepseek-chat) | 🟢 OpenAI (gpt-4o)", "🟠 Claude (claude-3.5-sonnet) | 🔴 Qwen (qwen-max) | ⚫ Ollama (llama3) | ⚙️ Custom"]
add_bullet(slide, 1, 4.7, 11, 1.5, items, 15)
add_footer(slide, 26)

# ========== 第10页：苏格拉底方法 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "EduGuide中的苏格拉底方法")
add_text(slide, 1, 1.5, 5.5, 0.5, "传统 vs. EduGuide", 20, True, PURPLE)
items1 = ["传统：\"NLP代表自然语言处理\" → 直接给答案", "EduGuide：\"你认为N代表什么？\" → 引导发现", "EduGuide：答错 → 分析错因 → 给出引导 → 重试"]
add_bullet(slide, 1, 2.2, 5.5, 2, items1, 15)
add_text(slide, 7, 1.5, 5.5, 0.5, "教学流程", 20, True, PURPLE)
items2 = [
    "步骤1：提出引导问题 → 学生回答",
    "步骤2：AI验证答案",
    "  ✅ 正确 → 进入下一步",
    "  ❌ 错误 → 🔍分析错因 → 💡引导提示 → 🔄重试",
]
add_bullet(slide, 7, 2.2, 5.5, 2.5, items2, 15)
add_text(slide, 1, 5.5, 11, 0.8, "🧠 促进深度学习 | 📊 记忆保持率提高55% | 👨‍🏫 模仿人类教师辅导方法", 16, False, GRAY, PP_ALIGN.CENTER)
add_footer(slide, 26)

# ========== 第11页：AI答案验证（核心创新）==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "AI答案验证 + 错因分析（核心创新）")
add_text(slide, 1, 1.5, 11, 0.6, "这是EduGuide区别于所有竞品的核心功能", 18, True, RED)
flow = [
    "学生提交答案",
    "    ↓",
    "调用LLM验证（与当前AI提供商相同）",
    "    ↓",
    "LLM返回: { correct, feedback, error_analysis, guidance }",
    "    ↓",
    "✅ 正确 → 进入下一步    ❌ 错误 → 显示错因分析 + 引导提示 → 无限重试",
]
add_bullet(slide, 1, 2.3, 6, 3.5, flow, 15)
add_text(slide, 7.5, 2.3, 5, 0.5, "错因分析示例", 18, True, ORANGE)
items = [
    "❌ 学生: \"语言模型是翻译工具\"",
    "🔍 错因: 学生将用途局限在翻译。语言模型核心是基于概率预测文本序列。",
    "💡 引导: \"想想手机输入法预测下一个词，这和翻译有关吗？\"",
]
add_bullet(slide, 7.5, 3, 5.5, 2.5, items, 14)
add_text(slide, 1, 5.8, 11, 0.6, "ChatGPT：直接给答案 ❌ | Quizlet：只判断对错 ❌ | EduGuide：验证+分析+引导 ✅", 16, True, PURPLE, PP_ALIGN.CENTER)
add_footer(slide, 26)

# ========== 第12页：功能总览 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "核心功能总览")
items1 = [
    "📁 多文件上传 (TXT/PDF/DOCX/PPTX)",
    "📝 直接粘贴文本",
    "🚀 一键生成 (~13秒)",
    "📊 分层题目 (基础/中级/高级)",
    "✅ AI答案验证",
    "🔍 错因分析",
]
items2 = [
    "💡 引导提示 (不直接给答案)",
    "🔄 无限次重试",
    "🌐 多语言 (英/中/繁)",
    "🔧 7个API提供商",
    "📡 Agent监控可视化",
    "🎨 现代紫色渐变UI",
]
add_bullet(slide, 1, 1.6, 5.5, 5, items1, 16)
add_bullet(slide, 7, 1.6, 5.5, 5, items2, 16)
add_footer(slide, 26)

# ========== 第13页：演示准备 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "现场演示")
items = [
    "1. 上传教学材料（支持多文件）",
    "2. 自动生成题目（约13秒）",
    "3. 交互式练习 + 故意答错 → 展示错因分析和引导",
    "4. 语言切换演示",
    "5. Agent监控页面",
]
add_bullet(slide, 1, 2, 5, 3.5, items, 20)
add_text(slide, 7, 2, 5, 1, "🔗 http://43.160.215.90:8080\n🔗 https://github.com/SheridanLu/EduGuide-QuestionGen", 16, False, PURPLE)
add_footer(slide, 26)

# ========== 第14页：演示执行 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "现场演示 - 执行脚本")
add_text(slide, 1, 1.5, 5.5, 0.5, "1. 上传+生成 (1分钟)", 18, True, PURPLE)
add_bullet(slide, 1, 2.2, 5.5, 1.5, ["上传材料 → 点击生成 → 等待13秒 → 展示知识点和题目"], 15)
add_text(slide, 1, 3.5, 5.5, 0.5, "2. 交互式练习 (1.5分钟) ⭐重点", 18, True, ORANGE)
add_bullet(slide, 1, 4.2, 5.5, 2, [
    "开始解题 → 正确答案 → 通过",
    "故意答错 → 🎯展示错因分析+引导提示 → 重试 → 通过",
], 15)
add_text(slide, 7, 1.5, 5.5, 0.5, "3. 语言切换+监控 (1分钟)", 18, True, PURPLE)
add_bullet(slide, 7, 2.2, 5.5, 2, [
    "切换英语/繁体中文",
    "打开Agent Monitor → 展示4个Agent工作流",
    "多文件上传演示",
], 15)
add_footer(slide, 26)

# ========== 第15页：演示亮点 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "演示亮点")
items = [
    "✅ 自动生成 — 13秒完成",
    "✅ 苏格拉底式引导 — 不直接给答案",
    "✅ AI答案验证 — 实时判断",
    "✅ 错因分析 — 分析为什么错",
    "✅ 引导提示 — 告诉你如何思考",
    "✅ 无限重试 — 不放弃任何学生",
    "✅ 多语言切换 — 英/中/繁",
    "✅ Agent监控 — 可视化工作流",
]
add_bullet(slide, 1, 1.6, 11, 5, items, 18)
add_text(slide, 1, 6, 11, 0.8, "教师需要数小时的工作，EduGuide在数秒内完成，并且像好老师一样引导学生自己发现答案。", 17, True, PURPLE, PP_ALIGN.CENTER)
add_footer(slide, 26)

# ========== 第16页：优势 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "优势")
add_text(slide, 1, 1.5, 3.5, 0.5, "✅ 有效性", 18, True, GREEN)
add_bullet(slide, 1, 2.2, 3.5, 2, ["⏰ 节省时间：2-3h→13s", "🎯 3级难度适配", "🧠 深度学习", "🔍 AI验证+错因分析"], 14)
add_text(slide, 5, 1.5, 3.5, 0.5, "✅ 技术可行性", 18, True, GREEN)
add_bullet(slide, 5, 2.2, 3.5, 2, ["🤖 成熟LLM", "🏗️ 4个独立Agent", "📁 4种文件格式", "🔌 7个API提供商"], 14)
add_text(slide, 9, 1.5, 3.5, 0.5, "✅ 可用性", 18, True, GREEN)
add_bullet(slide, 9, 2.2, 3.5, 2, ["🖥️ 现代UI", "📱 无需安装", "⚡ 实时反馈", "📡 Agent监控"], 14)
add_footer(slide, 26)

# ========== 第17页：劣势 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "劣势与局限")
add_text(slide, 1, 1.5, 3.5, 0.5, "⚠️ 技术局限", 18, True, RED)
add_bullet(slide, 1, 2.2, 3.5, 2, ["依赖LLM质量", "需API密钥(成本)", "需互联网", "仅限文本材料"], 14)
add_text(slide, 5, 1.5, 3.5, 0.5, "⚠️ 教育局限", 18, True, ORANGE)
add_bullet(slide, 5, 2.2, 3.5, 2, ["不能替代教师", "LLM可能偏见", "仅3种语言", "适合文本类学科"], 14)
add_text(slide, 9, 1.5, 3.5, 0.5, "缓解策略", 18, True, GREEN)
add_bullet(slide, 9, 2.2, 3.5, 2, ["多提供商容错", "教师可审查", "Ollama离线", "开源改进"], 14)
add_footer(slide, 26)

# ========== 第18页：对比 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "与现有方案对比")
rows = [
    ("功能", "EduGuide", "ChatGPT", "Quizlet"),
    ("自动题目生成", "✅", "✅", "❌"),
    ("从材料生成", "✅", "✅", "❌"),
    ("苏格拉底引导", "✅", "⚠️", "❌"),
    ("AI答案验证", "✅", "❌", "❌"),
    ("错因分析", "✅", "❌", "❌"),
    ("引导提示", "✅", "❌", "❌"),
    ("无限重试", "✅", "❌", "❌"),
    ("Agent监控", "✅", "❌", "❌"),
    ("开源", "✅", "❌", "❌"),
]
for i, row in enumerate(rows):
    y = 1.5 + i * 0.55
    is_header = (i == 0)
    color = PURPLE if is_header else BLACK
    bold = is_header
    size = 15 if is_header else 14
    for j, cell in enumerate(row):
        add_text(slide, 1 + j * 2.8, y, 2.8, 0.5, cell, size, bold, color)
add_footer(slide, 26)

# ========== 第19页：总结 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "总结")
add_text(slide, 1, 1.5, 11, 0.6, "我们构建了什么：", 20, True, PURPLE)
items = ["📚 自动从材料生成分层题目", "🎯 苏格拉底式引导教学", "🤖 AI答案验证+错因分析+引导提示", "🔄 无限次重试，不放弃任何学生", "📡 Agent监控可视化", "🌍 多语言支持 | 🔧 7个AI提供商"]
add_bullet(slide, 1, 2.3, 11, 3, items, 17)
add_text(slide, 1, 5.5, 11, 0.6, "核心创新：多Agent架构 + 苏格拉底方法 + AI答案验证(错因分析) + 无限重试", 16, True, PURPLE, PP_ALIGN.CENTER)
add_footer(slide, 26)

# ========== 第20页：关键收益 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "关键收益")
items1 = ["👨‍🏫 面向教师", "⏰ 节省数小时出题时间", "🎯 即时生成分层题目", "📝 专注于教学"]
items2 = ["👨‍🎓 面向学生", "🧠 引导式发现学习", "🔍 犯错后获得错因分析", "🔄 无限重试，按自己节奏"]
items3 = ["🏫 面向学校", "💰 7个API提供商可选", "📱 无需安装", "📚 适用于现有材料"]
items4 = ["🌍 面向开发者", "💻 开源透明", "🏗️ 模块化架构", "📡 Agent Monitor调试"]
for idx, items in enumerate([items1, items2, items3, items4]):
    x = 0.8 + (idx % 2) * 6.2
    y = 1.6 + (idx // 2) * 3
    add_bullet(slide, x, y, 5.5, 2.5, items, 15)
add_footer(slide, 26)

# ========== 第21页：未来工作 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "未来工作")
add_text(slide, 1, 1.5, 5.5, 0.5, "短期（3个月）", 18, True, PURPLE)
add_bullet(slide, 1, 2.2, 5.5, 1.5, ["📊 学习分析仪表板", "🎨 多媒体支持", "📱 移动应用"], 15)
add_text(slide, 7, 1.5, 5.5, 0.5, "中期（6个月）", 18, True, PURPLE)
add_bullet(slide, 7, 2.2, 5.5, 1.5, ["🤖 教育领域微调模型", "👥 协作学习", "📈 自适应难度"], 15)
add_text(slide, 1, 4.2, 5.5, 0.5, "长期（1年）", 18, True, PURPLE)
add_bullet(slide, 1, 4.9, 5.5, 1.5, ["🔌 LMS集成 (Moodle/Canvas)", "🌐 更多语言", "🧪 教育研究平台"], 15)
add_footer(slide, 26)

# ========== 第22页：结论 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PURPLE)
add_text(slide, 1, 1.5, 11, 1, "结论", 36, True, WHITE, PP_ALIGN.CENTER)
items = ["1. 自动化出题任务（数小时→数秒）", "2. 保留好教学的精髓（引导而非告诉）", "3. 超越简单对错（错因分析+引导）", "4. 使优质教育全球可访问"]
add_bullet(slide, 2, 3, 9, 2.5, items, 20, WHITE)
add_text(slide, 1, 5.5, 11, 0.8, "\"好的老师不给答案。他们引导学生自己找到答案。EduGuide通过AI使这成为可能。\"", 18, False, RGBColor(220, 220, 255), PP_ALIGN.CENTER)
add_text(slide, 1, 6.5, 11, 0.5, "🌐 http://43.160.215.90:8080 | 💻 github.com/SheridanLu/EduGuide-QuestionGen", 14, False, RGBColor(200, 200, 255), PP_ALIGN.CENTER)

# ========== 第23页：致谢 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PURPLE)
add_text(slide, 1, 2, 11, 1.5, "谢谢！", 48, True, WHITE, PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.8, "❓ 有问题吗？", 24, False, RGBColor(220, 220, 255), PP_ALIGN.CENTER)
add_text(slide, 2, 5, 9, 1.5, "Q&A 准备:\nQ1: 为什么用多Agent? → 专业化、易维护、可监控\nQ2: AI验证准确率? → 精心prompt + 关键词fallback\nQ3: 错因分析如何实现? → LLM返回JSON格式的分析和引导", 14, False, RGBColor(200, 200, 255))

# ========== 第24页：参考文献 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "参考文献（APA第7版）")
refs = [
    "Brown, T., et al. (2020). Language models are few-shot learners. NeurIPS, 33, 1877-1901.",
    "OpenAI. (2023). GPT-4 technical report. arXiv:2303.08774.",
    "Du, Z., et al. (2022). GLM: General language model pretraining. ACL, 320-335.",
    "Kurdi, G., et al. (2020). A systematic review of automatic question generation. IJAIED, 30(1), 121-204.",
    "Graesser, A. C., et al. (2001). Intelligent tutoring systems. AI Magazine, 22(4), 39-51.",
    "Paul, R., & Elder, L. (2007). Critical thinking: The art of Socratic questioning. JDE, 31(1), 36-37.",
    "Freeman, S., et al. (2014). Active learning increases student performance. PNAS, 111(23), 8410-8415.",
    "Wooldridge, M. (2009). An introduction to multiagent systems. Wiley.",
]
add_bullet(slide, 0.8, 1.5, 11.5, 5.5, refs, 13)
add_footer(slide, 26)

# ========== 第25页：工作分配 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "附录：工作分配")
rows = [
    ("成员", "角色", "幻灯片", "贡献"),
    ("[姓名1]", "引言", "1-6", "问题研究、背景分析"),
    ("[姓名2]", "系统设计", "7-11", "架构设计、技术栈"),
    ("[姓名3]", "实现", "12-15", "开发、演示"),
    ("[姓名4]", "评估", "16-18", "分析、对比"),
    ("[姓名5]", "结论", "19-24", "总结、Q&A"),
]
for i, row in enumerate(rows):
    y = 1.5 + i * 0.55
    is_header = (i == 0)
    color = PURPLE if is_header else BLACK
    for j, cell in enumerate(row):
        add_text(slide, 0.8 + j * 3, y, 3, 0.5, cell, 15 if is_header else 14, is_header, color)
add_footer(slide, 26)

# ========== 保存 ==========
output_path = "/root/.openclaw/workspace/projects/EduGuide-QuestionGen/EduGuide_Presentation_Chinese_v2.pptx"
prs.save(output_path)
print(f"✅ 中文PPT v2.0 已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 张幻灯片")
